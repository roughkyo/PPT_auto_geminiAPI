"""
학술 스타일 PowerPoint 프레젠테이션 자동 생성 스크립트
slides.json 파일을 읽어 이미지가 포함된 PPT를 생성합니다.
Google Gemini API를 활용하여 고품질 콘텐츠를 자동 생성합니다.
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
import google.generativeai as genai
from dotenv import load_dotenv

# 환경 변수 로드
load_dotenv()

def initialize_gemini_api():
    """Gemini API를 초기화합니다."""
    api_key = os.getenv('GEMINI_API_KEY')
    if not api_key:
        print("⚠ 경고: GEMINI_API_KEY 환경 변수가 설정되지 않았습니다.")
        print("   Gemini API 기능을 사용하려면 .env 파일에 API 키를 설정하세요.")
        return None
    
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-pro')
        print("✓ Gemini API 초기화 완료")
        return model
    except Exception as e:
        print(f"❌ Gemini API 초기화 실패: {e}")
        return None


def generate_slides_with_gemini(topic, num_slides=5, model=None):
    """Gemini API를 사용하여 주제에 맞는 슬라이드 콘텐츠를 생성합니다."""
    if not model:
        print("⚠ Gemini API가 초기화되지 않았습니다. 기본 모드로 진행합니다.")
        return None
    
    prompt = f"""
주제: {topic}

위 주제에 대한 트렌디하고 고퀄리티 프레젠테이션을 위한 {num_slides}개의 슬라이드 콘텐츠를 생성해주세요.

각 슬라이드는 다음 형식의 JSON으로 작성해주세요:

{{
  "topic": "{topic}",
  "design_theme": {{
    "primary_color": "#667eea",
    "secondary_color": "#764ba2",
    "accent_color": "#f093fb",
    "style": "glassmorphism"
  }},
  "slides": [
    {{
      "title": "슬라이드 제목",
      "content": [
        "**핵심 개념**: 설명과 함께",
        "재미있는 비유: 마치 ~처럼",
        "구체적인 예시와 수치",
        "**강조할 포인트**: 중요한 내용"
      ],
      "image_prompt": "modern glassmorphism style, gradient background with purple and blue tones, semi-transparent frosted glass elements, subtle blur effects, [구체적인 다이어그램 설명], professional tech illustration, vibrant neon accents, clean minimalist design, soft shadows, depth layers"
    }}
  ]
}}

요구사항:
1. 총 {num_slides}장의 슬라이드 (논리적 구조: 도입 → 핵심 개념 → 심화 → 응용 → 미래 전망)
2. 각 슬라이드는 4-6개의 핵심 포인트로 구성
3. **중요 개념**은 마크다운 굵은 글씨로 표현 (예: **트랜스포머**, **어텐션 메커니즘**)
4. 학술적 정확성을 유지하면서도 위트있는 비유와 예시를 포함
   - 예: "마치 ~처럼", "~와 비슷하게", "쉽게 말하면 ~"
5. 각 슬라이드마다 글라스모피즘 스타일 이미지 프롬프트 생성
6. 이미지 프롬프트는 반드시 "modern glassmorphism style, gradient background with purple and blue tones..."로 시작
7. 색상 테마: 보라-파랑-핑크 그라데이션 (#667eea, #764ba2, #f093fb)
8. 전문적이면서도 흥미롭고 재미있는 톤 유지
9. 각 포인트는 간결하지만 정보가 풍부하게

JSON 형식만 반환하고, 다른 설명은 포함하지 마세요.
"""
    
    try:
        print(f"\n🤖 Gemini API로 '{topic}' 주제의 고퀄리티 슬라이드 생성 중...")
        print(f"   📊 슬라이드 개수: {num_slides}장")
        print(f"   🎨 스타일: 글라스모피즘 (보라-파랑 그라데이션)")
        print(f"   ✨ 특징: 학술적 + 위트있는 콘텐츠")
        
        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.8,  # 창의성을 높여 위트있는 콘텐츠 생성
                top_p=0.95,
                top_k=40,
                max_output_tokens=8192,  # 10장 슬라이드를 위해 토큰 수 증가
            )
        )
        
        # JSON 파싱
        content = response.text.strip()
        # 마크다운 코드 블록 제거
        if content.startswith('```'):
            content = content.split('```')[1]
            if content.startswith('json'):
                content = content[4:]
            content = content.strip()
        
        slides_data = json.loads(content)
        print(f"✓ Gemini API로 {len(slides_data.get('slides', []))}개 슬라이드 생성 완료")
        print(f"✓ 디자인 테마: {slides_data.get('design_theme', {}).get('style', 'default')}")
        return slides_data
        
    except json.JSONDecodeError as e:
        print(f"❌ JSON 파싱 오류: {e}")
        print(f"응답 내용: {response.text[:500]}...")
        return None
    except Exception as e:
        print(f"❌ 슬라이드 생성 실패: {e}")
        return None



def enhance_slide_content_with_gemini(slide_data, model=None):
    """기존 슬라이드 콘텐츠를 Gemini API로 개선합니다."""
    if not model:
        return slide_data
    
    prompt = f"""
다음 슬라이드 콘텐츠를 더 전문적이고 학술적으로 개선해주세요:

제목: {slide_data['title']}
콘텐츠:
{chr(10).join(f"- {point}" for point in slide_data.get('content', []))}

요구사항:
1. 제목을 더 명확하고 전문적으로 개선
2. 각 포인트를 더 구체적이고 정보가 풍부하게 작성
3. 학술적 톤 유지
4. 3-5개의 핵심 포인트로 정리
5. 간결하면서도 정보가 풍부하게

다음 JSON 형식으로만 응답해주세요:
{{
  "title": "개선된 제목",
  "content": [
    "개선된 포인트 1",
    "개선된 포인트 2",
    "개선된 포인트 3"
  ]
}}
"""
    
    try:
        response = model.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.5,  # 더 일관성 있는 개선
                top_p=0.8,
                max_output_tokens=1024,
            )
        )
        
        content = response.text.strip()
        if content.startswith('```'):
            content = content.split('```')[1]
            if content.startswith('json'):
                content = content[4:]
            content = content.strip()
        
        enhanced = json.loads(content)
        return {**slide_data, **enhanced}
        
    except Exception as e:
        print(f"  ⚠ 콘텐츠 개선 실패: {e}")
        return slide_data


def load_slides_data(json_path='slides.json'):
    """JSON 파일에서 슬라이드 데이터를 로드합니다."""
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        print(f"✓ JSON 파일 로드 완료: {len(data.get('slides', []))}개 슬라이드")
        return data
    except FileNotFoundError:
        print(f"❌ 오류: {json_path} 파일을 찾을 수 없습니다.")
        return None
    except json.JSONDecodeError:
        print(f"❌ 오류: {json_path} 파일의 JSON 형식이 올바르지 않습니다.")
        return None



def create_title_slide(prs, topic):
    """타이틀 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 제목 추가
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = topic
    
    # 제목 스타일링
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)  # 다크 블루
    
    # 부제목 추가
    subtitle_top = Inches(4.2)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"생성일: {datetime.now().strftime('%Y년 %m월 %d일')}"
    
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_paragraph.font.size = Pt(16)
    subtitle_paragraph.font.color.rgb = RGBColor(100, 100, 100)
    
    print("✓ 타이틀 슬라이드 생성 완료")


def create_content_slide(prs, slide_data, slide_number, images_dir='images'):
    """콘텐츠 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 제목 추가
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(9)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 51, 102)
    
    # 이미지 추가 (왼쪽)
    image_path = os.path.join(images_dir, f'slide_{slide_number}.png')
    if os.path.exists(image_path):
        img_left = Inches(0.5)
        img_top = Inches(1.5)
        img_width = Inches(4.5)
        
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            print(f"  ✓ 이미지 추가: {image_path}")
        except Exception as e:
            print(f"  ⚠ 이미지 추가 실패: {e}")
    else:
        print(f"  ⚠ 이미지 파일 없음: {image_path}")
    
    # 콘텐츠 텍스트 추가 (오른쪽)
    content_left = Inches(5.2)
    content_top = Inches(1.5)
    content_width = Inches(4.3)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # 콘텐츠 포인트 추가
    for i, point in enumerate(slide_data.get('content', [])):
        if i > 0:
            content_frame.add_paragraph()
        
        p = content_frame.paragraphs[i]
        p.text = f"• {point}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(12)
        p.level = 0
    
    print(f"✓ 슬라이드 {slide_number} 생성 완료: {slide_data['title']}")


def generate_presentation(slides_data, output_dir='output'):
    """전체 프레젠테이션을 생성합니다."""
    # 출력 디렉토리 생성
    Path(output_dir).mkdir(exist_ok=True)
    
    # 새 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    topic = slides_data.get('topic', '프레젠테이션')
    
    # 타이틀 슬라이드 생성
    create_title_slide(prs, topic)
    
    # 콘텐츠 슬라이드 생성
    slides = slides_data.get('slides', [])
    for i, slide_data in enumerate(slides, 1):
        create_content_slide(prs, slide_data, i)
    
    # 파일 저장
    safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_', '-')).strip()
    safe_topic = safe_topic.replace(' ', '_')
    output_path = os.path.join(output_dir, f'{safe_topic}_presentation.pptx')
    
    prs.save(output_path)
    print(f"\n{'='*60}")
    print(f"✅ PPT 생성 완료!")
    print(f"📁 파일 위치: {output_path}")
    print(f"📊 총 슬라이드 수: {len(slides) + 1} (타이틀 포함)")
    print(f"{'='*60}\n")
    
    return output_path


def main():
    """메인 실행 함수"""
    print("\n" + "="*60)
    print("🎓 학술 스타일 PPT 자동 생성 시작")
    print("="*60 + "\n")
    
    # Gemini API 초기화
    gemini_model = initialize_gemini_api()
    
    # 사용자 입력 받기
    print("\n📋 PPT 생성 모드를 선택하세요:")
    print("1. 기존 slides.json 파일 사용")
    print("2. Gemini API로 새로운 슬라이드 생성")
    print("3. 기존 JSON 파일의 콘텐츠를 Gemini API로 개선")
    
    mode = input("\n선택 (1/2/3, 기본값: 1): ").strip() or "1"
    
    slides_data = None
    
    if mode == "2":
        # Gemini API로 새로운 슬라이드 생성
        if not gemini_model:
            print("❌ Gemini API를 사용할 수 없습니다. 모드 1을 사용하세요.")
            return
        
        topic = input("\n📝 프레젠테이션 주제를 입력하세요: ").strip()
        if not topic:
            print("❌ 주제를 입력해야 합니다.")
            return
        
        num_slides = input("📊 생성할 슬라이드 개수 (기본값: 10): ").strip() or "10"
        try:
            num_slides = int(num_slides)
        except ValueError:
            num_slides = 10
        
        slides_data = generate_slides_with_gemini(topic, num_slides, gemini_model)
        
        if slides_data:
            # 생성된 데이터를 파일로 저장
            output_json = f"slides_generated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
            with open(output_json, 'w', encoding='utf-8') as f:
                json.dump(slides_data, f, ensure_ascii=False, indent=2)
            print(f"✓ 생성된 슬라이드 데이터 저장: {output_json}")
    
    elif mode == "3":
        # 기존 콘텐츠를 Gemini API로 개선
        if not gemini_model:
            print("❌ Gemini API를 사용할 수 없습니다. 모드 1을 사용하세요.")
            return
        
        json_path = input("\n📁 JSON 파일 경로 (기본값: slides.json): ").strip() or "slides.json"
        slides_data = load_slides_data(json_path)
        
        if slides_data:
            print("\n🔧 Gemini API로 콘텐츠 개선 중...")
            for i, slide in enumerate(slides_data.get('slides', []), 1):
                print(f"  슬라이드 {i} 개선 중...")
                slides_data['slides'][i-1] = enhance_slide_content_with_gemini(slide, gemini_model)
            
            # 개선된 데이터를 파일로 저장
            output_json = f"slides_enhanced_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
            with open(output_json, 'w', encoding='utf-8') as f:
                json.dump(slides_data, f, ensure_ascii=False, indent=2)
            print(f"✓ 개선된 슬라이드 데이터 저장: {output_json}")
    
    else:
        # 기존 JSON 파일 사용
        json_path = input("\n📁 JSON 파일 경로 (기본값: slides.json): ").strip() or "slides.json"
        slides_data = load_slides_data(json_path)
    
    if not slides_data:
        return
    
    # 이미지 디렉토리 확인
    images_dir = 'images'
    if not os.path.exists(images_dir):
        print(f"⚠ 경고: {images_dir} 디렉토리가 없습니다. 이미지 없이 진행합니다.")
        Path(images_dir).mkdir(exist_ok=True)
    
    # PPT 생성
    output_path = generate_presentation(slides_data)
    
    print("✨ 모든 작업이 완료되었습니다!")


if __name__ == '__main__':
    main()

