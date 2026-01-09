"""
자동으로 slides.json을 사용하여 PPT를 생성하는 스크립트
"""

import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime


def load_slides_data(json_path='slides.json'):
    """JSON 파일에서 슬라이드 데이터를 로드합니다."""
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        print(f"✓ JSON 파일 로드 완료: {len(data.get('slides', []))}개 슬라이드")
        return data
    except FileNotFoundError:
        print(f"❌ 오류: {json_path} 파일을 찾을 수 없습니다.")
        return None
    except json.JSONDecodeError:
        print(f"❌ 오류: {json_path} 파일의 JSON 형식이 올바르지 않습니다.")
        return None


def create_title_slide(prs, topic):
    """타이틀 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 제목 추가
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = topic
    
    # 제목 스타일링
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(102, 126, 234)  # #667eea
    
    # 부제목 추가
    subtitle_top = Inches(4.2)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"생성일: {datetime.now().strftime('%Y년 %m월 %d일')}"
    
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_paragraph.font.size = Pt(16)
    subtitle_paragraph.font.color.rgb = RGBColor(118, 75, 162)  # #764ba2
    
    print("✓ 타이틀 슬라이드 생성 완료")


def create_content_slide(prs, slide_data, slide_number, images_dir='images'):
    """콘텐츠 슬라이드를 생성합니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    
    # 제목 추가
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(9)
    title_height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(32)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(102, 126, 234)  # #667eea
    
    # 이미지 추가 (왼쪽)
    image_path = os.path.join(images_dir, f'slide_{slide_number}.png')
    if os.path.exists(image_path):
        img_left = Inches(0.5)
        img_top = Inches(1.5)
        img_width = Inches(4.5)
        
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            print(f"  ✓ 이미지 추가: {image_path}")
        except Exception as e:
            print(f"  ⚠ 이미지 추가 실패: {e}")
    else:
        print(f"  ⚠ 이미지 파일 없음: {image_path}")
    
    # 콘텐츠 텍스트 추가 (오른쪽)
    content_left = Inches(5.2)
    content_top = Inches(1.5)
    content_width = Inches(4.3)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # 콘텐츠 포인트 추가
    for i, point in enumerate(slide_data.get('content', [])):
        if i > 0:
            content_frame.add_paragraph()
        
        p = content_frame.paragraphs[i]
        p.text = f"• {point}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(10)
        p.level = 0
    
    print(f"✓ 슬라이드 {slide_number} 생성 완료: {slide_data['title']}")


def generate_presentation(slides_data, output_dir='output'):
    """전체 프레젠테이션을 생성합니다."""
    # 출력 디렉토리 생성
    Path(output_dir).mkdir(exist_ok=True)
    
    # 새 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    topic = slides_data.get('topic', '프레젠테이션')
    
    # 타이틀 슬라이드 생성
    create_title_slide(prs, topic)
    
    # 콘텐츠 슬라이드 생성
    slides = slides_data.get('slides', [])
    for i, slide_data in enumerate(slides, 1):
        create_content_slide(prs, slide_data, i)
    
    # 파일 저장
    safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '_', '-')).strip()
    safe_topic = safe_topic.replace(' ', '_')
    output_path = os.path.join(output_dir, f'{safe_topic}_presentation.pptx')
    
    prs.save(output_path)
    print(f"\n{'='*60}")
    print(f"✅ PPT 생성 완료!")
    print(f"📁 파일 위치: {output_path}")
    print(f"📊 총 슬라이드 수: {len(slides) + 1} (타이틀 포함)")
    print(f"🎨 디자인 테마: {slides_data.get('design_theme', {}).get('style', 'default')}")
    print(f"{'='*60}\n")
    
    return output_path


def main():
    """메인 실행 함수"""
    print("\n" + "="*60)
    print("🎨 트렌디 학술 스타일 PPT 자동 생성")
    print("="*60 + "\n")
    
    # JSON 데이터 로드
    slides_data = load_slides_data()
    if not slides_data:
        return
    
    # 이미지 디렉토리 확인
    images_dir = 'images'
    if not os.path.exists(images_dir):
        print(f"⚠ 경고: {images_dir} 디렉토리가 없습니다. 이미지 없이 진행합니다.")
        Path(images_dir).mkdir(exist_ok=True)
    
    # PPT 생성
    output_path = generate_presentation(slides_data)
    
    print("✨ 모든 작업이 완료되었습니다!")


if __name__ == '__main__':
    main()
